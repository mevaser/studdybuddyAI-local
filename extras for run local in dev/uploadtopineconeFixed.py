import os
import openai
import uuid
import tiktoken

from PyPDF2 import PdfReader
from docx import Document
from pptx import Presentation
from pinecone import Pinecone
from dotenv import load_dotenv
load_dotenv()

# ✅ הגדרות קבועות
PINECONE_API_KEY = os.getenv("PINECONE_API_KEY")
INDEX_NAME = "studybuddy-v2"
OPENAI_API_KEY = os.getenv("api_key")

openai.api_key = OPENAI_API_KEY
pc = Pinecone(api_key=PINECONE_API_KEY)
index = pc.Index(INDEX_NAME)

# ✅ פונקציית קריאת טקסט מקובץ
def extract_text_from_file(file_path):
    ext = file_path.split(".")[-1].lower()
    try:
        if ext == "pdf":
            reader = PdfReader(file_path)
            return "\n".join(page.extract_text() for page in reader.pages if page.extract_text())
        elif ext == "docx":
            doc = Document(file_path)
            return "\n".join(p.text for p in doc.paragraphs if p.text.strip())
        elif ext == "pptx":
            prs = Presentation(file_path)
            return "\n".join(shape.text for slide in prs.slides for shape in slide.shapes if hasattr(shape, "text"))
        else:
            return ""
    except Exception as e:
        print(f"⚠ שגיאה בקריאת הקובץ {file_path}: {e}")
        return ""

# ✅ פיצול טקסט לצ'אנקים לפי טוקנים + חפיפה
def split_text_by_tokens(text, max_tokens=400, overlap_tokens=100):
    encoder = tiktoken.encoding_for_model("text-embedding-ada-002")
    tokens = encoder.encode(text)

    chunks = []
    start = 0
    while start < len(tokens):
        end = start + max_tokens
        chunk_tokens = tokens[start:end]
        chunk_text = encoder.decode(chunk_tokens)
        chunks.append(chunk_text)
        start += max_tokens - overlap_tokens  # חפיפה

    return chunks

# ✅ המרת טקסט לווקטור בעזרת GPT
def get_embedding(text):
    response = openai.Embedding.create(
        input=text,
        model="text-embedding-ada-002"
    )
    return response['data'][0]['embedding']

# ✅ שליחת צ'אנקים ל-Pinecone כולל בדיקה מראש
def upload_chunks_to_pinecone(chunks, metadata):
    uploaded = 0
    skipped = 0
    for chunk in chunks:
        chunk_id = str(uuid.uuid5(uuid.NAMESPACE_DNS, chunk))
        existing = index.fetch(ids=[chunk_id], namespace="default").vectors
        if chunk_id in existing:
            print(f"✔ צ'אנק כבר קיים (ID: {chunk_id}), מדלג.")
            skipped += 1
            continue

        vec = get_embedding(chunk)
        index.upsert(
            vectors=[
                {
                    "id": chunk_id,
                    "values": vec,
                    "metadata": {**metadata, "text": chunk}
                }
            ],
            namespace="default"
        )
        print(f"⬆ צ'אנק הועלה (ID: {chunk_id})")
        uploaded += 1

    print(f"סיכום לקובץ: {uploaded} צ'אנקים הועלו, {skipped} דולגו.")

# ✅ ריצה על תיקייה עם חיווי מלא
def process_folder(folder_path, course_name):
    files = [f for f in os.listdir(folder_path) if os.path.isfile(os.path.join(folder_path, f)) and not f.startswith("~$")]
    total_files = len(files)
    for idx, filename in enumerate(files, start=1):
        file_path = os.path.join(folder_path, filename)
        print(f"\n▶ מעבד קובץ {idx}/{total_files}: {filename}")
        text = extract_text_from_file(file_path)
        chunks = split_text_by_tokens(text)
        if chunks:
            upload_chunks_to_pinecone(chunks, metadata={"course": course_name, "source": filename})
        else:
            print(f"⚠ אין תוכן להעלות מתוך {filename}")

    print("\n✅ סיום תהליך - כל הקבצים עובדו.")

# ✅ הרצה לדוגמה
if __name__ == "__main__":
    folder_path = "C:\\Users\\User\\Desktop\\אגם\\לימודים\\שנה ג\\פרוייקט\\networks course"
    process_folder(folder_path, course_name="networks")
    print("✅ Upload complete.")
