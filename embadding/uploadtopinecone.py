# ✅ שלב 1: ייבוא ספריות
import os
import pinecone
import openai
import uuid

from PyPDF2 import PdfReader
from docx import Document
from pptx import Presentation

# ✅ שלב 2: הגדרות קבועות
PINECONE_API_KEY = os.getenv("PINECONE_API_KEY")
PINECONE_ENV = "us-east-1"
INDEX_NAME = "studybuddy-index"
OPENAI_API_KEY = os.getenv("api_key")

openai.api_key = OPENAI_API_KEY
pinecone.init(api_key=PINECONE_API_KEY, environment=PINECONE_ENV)
index = pinecone.Index(INDEX_NAME)

# ✅ שלב 3: פונקציה לטעינת קבצים

def extract_text_from_file(file_path):
    ext = file_path.split(".")[-1].lower()
    if ext == "pdf":
        reader = PdfReader(file_path)
        return "\n".join(page.extract_text() for page in reader.pages if page.extract_text())
    elif ext == "docx":
        doc = Document(file_path)
        return "\n".join(p.text for p in doc.paragraphs if p.text.strip())
    elif ext == "pptx":
        prs = Presentation(file_path)
        return "\n".join(shape.text for slide in prs.slides for shape in slide.shapes if hasattr(shape, "text"))
    else:
        return ""

# ✅ שלב 4: פיצול טקסט למקטעים קצרים (ל־embedding)

def split_text(text, max_length=500):
    paragraphs = text.split("\n")
    chunks = []
    current = ""
    for p in paragraphs:
        if len(current) + len(p) < max_length:
            current += p + " "
        else:
            chunks.append(current.strip())
            current = p + " "
    if current.strip():
        chunks.append(current.strip())
    return chunks

# ✅ שלב 5: המרת טקסט לוקטור בעזרת GPT

def get_embedding(text):
    response = openai.Embedding.create(
        input=text,
        model="text-embedding-ada-002"
    )
    return response['data'][0]['embedding']

# ✅ שלב 6: שליחת מקטעים ל-Pinecone

def upload_chunks_to_pinecone(chunks, metadata):
    vectors = []
    for chunk in chunks:
        vec = get_embedding(chunk)
        vectors.append((str(uuid.uuid4()), vec, {**metadata, "text": chunk}))

    index.upsert(vectors)

# ✅ שלב 7: ריצה ראשית על כל הקבצים מתיקיה

def process_folder(folder_path, course_name):
    for filename in os.listdir(folder_path):
        file_path = os.path.join(folder_path, filename)
        if os.path.isfile(file_path):
            print(f"Processing {filename}")
            text = extract_text_from_file(file_path)
            chunks = split_text(text)
            upload_chunks_to_pinecone(chunks, metadata={"course": course_name, "source": filename})

# ✅ שלב 8: דוגמה להרצה ראשית
if __name__ == "__main__":
    folder_path = "C:\Users\User\Desktop\אגם\לימודים\שנה ג\פרוייקט\networks course"  # שנה לנתיב שבו הקבצים אצלך
    process_folder(folder_path, course_name="networks")
    print("✅ Upload complete.")
